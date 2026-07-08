from __future__ import annotations

import json
import math
import os
import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path

from verificar_knowledge import DB as VERIFICAR_DB
from verificar_knowledge import answer_verificar_question


ROOT = Path(__file__).resolve().parents[1]
KNOWLEDGE_DIR = ROOT / "knowledge_base"
DATA_DIR = ROOT / "data"
INDEX_PATH = DATA_DIR / "index.json"
LINKS_DIR = KNOWLEDGE_DIR / "links"

IMAGE_FILE_TYPES = {".png", ".jpg", ".jpeg", ".webp"}
DOCUMENT_FILE_TYPES = {".txt", ".md", ".docx", ".xlsx", ".pptx", ".pdf"}
CODE_FILE_TYPES = {".js", ".json"}
SUPPORTED_FILE_TYPES = DOCUMENT_FILE_TYPES | CODE_FILE_TYPES | IMAGE_FILE_TYPES


class UserFacingError(RuntimeError):
    pass


@dataclass
class TextDocument:
    source: str
    text: str


def load_settings() -> None:
    try:
        from dotenv import load_dotenv
    except ModuleNotFoundError:
        return
    load_dotenv(ROOT / ".env", encoding="utf-8-sig")


def get_client():
    load_settings()
    if not os.getenv("OPENAI_API_KEY"):
        raise UserFacingError("Falta configurar OPENAI_API_KEY en el archivo .env.")

    try:
        from openai import OpenAI
    except ModuleNotFoundError as error:
        raise UserFacingError("Faltan dependencias. Ejecuta: pip install -r requirements.txt") from error

    return OpenAI()


def friendly_openai_error(error: Exception) -> UserFacingError:
    message = str(error)
    lowered = message.lower()
    if "insufficient_quota" in lowered or "exceeded your current quota" in lowered:
        return UserFacingError(
            "Tu API key esta bien, pero la cuenta no tiene cuota disponible. "
            "Revisa billing, saldo o limite mensual en OpenAI Platform."
        )
    if "invalid_api_key" in lowered or "incorrect api key" in lowered:
        return UserFacingError("La API key no es valida. Revisa el valor de OPENAI_API_KEY en .env.")
    if "model_not_found" in lowered:
        return UserFacingError("El modelo configurado no esta disponible para esta cuenta. Revisa OPENAI_MODEL en .env.")
    return UserFacingError(f"OpenAI devolvio un error: {error}")


def read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def read_docx(path: Path) -> str:
    from docx import Document as WordDocument

    document = WordDocument(path)
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    tables = []
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                tables.append(" | ".join(values))
    return "\n".join(paragraphs + tables)


def read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"Hoja: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def read_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"Diapositiva {index}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ModuleNotFoundError as error:
        raise UserFacingError("Falta pypdf para leer PDF. Ejecuta: pip install -r requirements.txt") from error

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"Pagina {index}\n{text.strip()}")
    return "\n\n".join(pages)


def read_image(path: Path) -> str:
    label = path.stem.replace("_", " ").replace("-", " ").strip()
    return (
        f"Imagen disponible en la base de conocimiento. "
        f"Nombre o color asociado: {label}. "
        f"Archivo: {path.name}."
    )


def read_web_page(url: str) -> str:
    import requests
    from bs4 import BeautifulSoup

    response = requests.get(url, timeout=20, headers={"User-Agent": "kb-rag-prototype/1.0"})
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text("\n")
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)


def load_local_documents() -> list[TextDocument]:
    # BLOQUE 1: carga de archivos locales por tipo.
    # Para agregar un formato nuevo, agrega su extension arriba y su lector aqui.
    documents = []
    for path in KNOWLEDGE_DIR.rglob("*"):
        if not path.is_file() or path.name == "urls.txt" or path.name.startswith("~$"):
            continue
        if path.suffix.lower() not in SUPPORTED_FILE_TYPES:
            continue

        if path.suffix.lower() in {".txt", ".md", ".js"}:
            text = read_text_file(path)
        elif path.suffix.lower() == ".docx":
            text = read_docx(path)
        elif path.suffix.lower() == ".xlsx":
            text = read_xlsx(path)
        elif path.suffix.lower() == ".pptx":
            text = read_pptx(path)
        elif path.suffix.lower() == ".pdf":
            text = read_pdf(path)
        elif path.suffix.lower() in IMAGE_FILE_TYPES:
            text = read_image(path)
        else:
            text = ""

        if text.strip():
            documents.append(TextDocument(source=str(path.relative_to(ROOT)), text=text))
    return documents


def load_web_documents() -> list[TextDocument]:
    url_file = KNOWLEDGE_DIR / "urls.txt"
    if not url_file.exists():
        return []

    documents = []
    urls = [
        line.strip()
        for line in url_file.read_text(encoding="utf-8", errors="ignore").splitlines()
        if line.strip() and not line.strip().startswith("#")
    ]
    for url in urls:
        text = read_web_page(url)
        if text.strip():
            documents.append(TextDocument(source=url, text=text))
    return documents


def load_resource_links() -> list[dict]:
    # BLOQUE 2: base liviana de enlaces.
    # Edita knowledge_base/links/*.json para agregar imagenes, videos, presentaciones o documentos por URL.
    resources = []
    if not LINKS_DIR.exists():
        return resources

    for path in sorted(LINKS_DIR.glob("*.json")):
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            continue
        items = payload.get("resources", payload if isinstance(payload, list) else [])
        if not isinstance(items, list):
            continue
        for item in items:
            if not isinstance(item, dict):
                continue
            url = str(item.get("url", "")).strip()
            if not re.match(r"^https?://", url, flags=re.IGNORECASE):
                continue
            resources.append(
                {
                    "source": str(path.relative_to(ROOT)),
                    "kind": str(item.get("kind") or payload.get("kind") or path.stem).strip(),
                    "title": str(item.get("title", "")).strip(),
                    "description": str(item.get("description", "")).strip(),
                    "url": url,
                    "aliases": item.get("aliases", []),
                    "code": str(item.get("code", "")).strip(),
                }
            )
    return resources


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 180) -> list[str]:
    normalized = " ".join(text.split())
    if not normalized:
        return []

    chunks = []
    start = 0
    while start < len(normalized):
        end = start + chunk_size
        chunks.append(normalized[start:end])
        next_start = end - overlap
        start = next_start if next_start > start else end
    return chunks


def embed_texts(client, texts: list[str]) -> list[list[float]]:
    model = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")
    embeddings = []
    batch_size = 80
    for start in range(0, len(texts), batch_size):
        batch = texts[start : start + batch_size]
        try:
            response = client.embeddings.create(model=model, input=batch)
        except Exception as error:
            raise friendly_openai_error(error) from error
        embeddings.extend(item.embedding for item in response.data)
    return embeddings


def build_index() -> dict:
    client = get_client()
    documents = load_local_documents() + load_web_documents()
    chunks = []
    for document in documents:
        for index, chunk in enumerate(chunk_text(document.text), start=1):
            chunks.append({"source": document.source, "chunk": index, "text": chunk})

    if not chunks:
        raise RuntimeError("No encontre documentos para indexar en knowledge_base.")

    embeddings = embed_texts(client, [chunk["text"] for chunk in chunks])
    for chunk, embedding in zip(chunks, embeddings):
        chunk["embedding"] = embedding

    DATA_DIR.mkdir(exist_ok=True)
    payload = {"chunks": chunks}
    INDEX_PATH.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def load_index() -> dict:
    if not INDEX_PATH.exists():
        raise UserFacingError("Primero crea el indice con el boton 'Crear indice'.")
    return json.loads(INDEX_PATH.read_text(encoding="utf-8"))


def cosine_similarity(left: list[float], right: list[float]) -> float:
    numerator = sum(a * b for a, b in zip(left, right))
    left_norm = math.sqrt(sum(a * a for a in left))
    right_norm = math.sqrt(sum(b * b for b in right))
    if not left_norm or not right_norm:
        return 0.0
    return numerator / (left_norm * right_norm)


def retrieve(question: str, top_k: int = 5) -> list[dict]:
    index = load_index()
    client = get_client()
    question_embedding = embed_texts(client, [question])[0]
    scored = []
    for chunk in index["chunks"]:
        score = cosine_similarity(question_embedding, chunk["embedding"])
        scored.append({**chunk, "score": score})
    return sorted(scored, key=lambda item: item["score"], reverse=True)[:top_k]


def build_context(matches: list[dict]) -> str:
    return "\n\n".join(
        f"[Fuente: {match['source']} | fragmento {match['chunk']}]\n{match['text']}"
        for match in matches
    )


def strip_visible_sources(text: str) -> str:
    cleaned = re.sub(r"\n+\*\*Fuentes:\*\*.*$", "", text, flags=re.IGNORECASE | re.DOTALL)
    cleaned = re.sub(r"\n+Fuentes:\s*.*$", "", cleaned, flags=re.IGNORECASE | re.DOTALL)
    return cleaned.strip()


def normalize_search_text(text: str) -> str:
    normalized = unicodedata.normalize("NFD", str(text or "").lower())
    normalized = "".join(char for char in normalized if unicodedata.category(char) != "Mn")
    return re.sub(r"[^a-z0-9]+", " ", normalized)


def search_tokens(text: str) -> set[str]:
    ignored = {
        "aqui",
        "color",
        "como",
        "con",
        "cual",
        "cuales",
        "dame",
        "de",
        "del",
        "el",
        "en",
        "es",
        "esta",
        "este",
        "la",
        "las",
        "lo",
        "los",
        "me",
        "muestra",
        "muestrame",
        "para",
        "que",
        "quiero",
        "un",
        "una",
        "ver",
    }
    return {
        token
        for token in normalize_search_text(text).split()
        if len(token) > 2 and token not in ignored
    }


def history_text(history: list[dict] | None, limit: int = 4) -> str:
    if not history:
        return ""
    parts = []
    for item in history[-limit:]:
        question = str(item.get("question", "")).strip()
        answer = str(item.get("answer", "")).strip()
        if question:
            parts.append(f"Pregunta anterior: {question}")
        if answer:
            parts.append(f"Respuesta anterior: {answer[:500]}")
    return "\n".join(parts)


def contextual_question(question: str, history: list[dict] | None = None) -> str:
    # BLOQUE 3: continuidad de conversacion.
    # Solo une historial cuando la pregunta actual parece seguimiento del mismo tema.
    current = str(question or "").strip()
    context = history_text(history)
    if not context:
        return current
    normalized = normalize_search_text(current)
    followup_markers = {"ese", "esa", "eso", "este", "esta", "estos", "estas", "tambien", "igual", "mismo", "misma", "video", "imagen", "link", "enlace", "documento", "presentacion"}
    current_tokens = set(normalized.split())

    if not current_tokens.intersection(followup_markers):
        return current

    previous_tokens = search_tokens(context)
    meaningful_current = search_tokens(current) - followup_markers
    if meaningful_current and not meaningful_current.intersection(previous_tokens):
        return current

    if len(current_tokens) <= 7 or current_tokens.intersection(followup_markers):
        return f"{context}\nPregunta actual: {current}"
    return current


def retrieve_resource_links(question: str, top_k: int = 6) -> list[dict]:
    tokens = search_tokens(question)
    if not tokens:
        return []

    matches = []
    for resource in load_resource_links():
        aliases = resource.get("aliases", [])
        alias_text = " ".join(str(alias) for alias in aliases) if isinstance(aliases, list) else str(aliases)
        searchable = " ".join(
            [
                str(resource.get("kind", "")),
                str(resource.get("title", "")),
                str(resource.get("description", "")),
                str(resource.get("code", "")),
                alias_text,
                str(resource.get("url", "")),
            ]
        )
        normalized = normalize_search_text(searchable)
        score = sum(3 if token in normalize_search_text(str(resource.get("title", ""))) else 1 for token in tokens if token in normalized)
        if score:
            matches.append(
                {
                    "source": resource["source"],
                    "chunk": 1,
                    "text": f"{resource['title']} - {resource['url']}",
                    "score": float(score),
                    "kind": "resource_link",
                    "resource_kind": resource.get("kind", ""),
                    "title": resource.get("title", ""),
                    "description": resource.get("description", ""),
                    "url": resource.get("url", ""),
                    "code": resource.get("code", ""),
                }
            )
    return sorted(matches, key=lambda item: item["score"], reverse=True)[:top_k]


def wants_external_link(question: str) -> bool:
    tokens = set(normalize_search_text(question).split())
    return bool(
        tokens.intersection(
            {
                "enlace",
                "enlaces",
                "link",
                "links",
                "url",
                "video",
                "videos",
                "imagen",
                "imagenes",
                "foto",
                "fotos",
                "presentacion",
                "presentaciones",
                "powerpoint",
                "ppt",
                "pptx",
            }
        )
    )


def retrieve_local(question: str, top_k: int = 8) -> list[dict]:
    # BLOQUE 4: busqueda local rapida sin embeddings.
    # Prioriza enlaces organizados, luego DB de VERIFICAR, luego documentos locales.
    tokens = search_tokens(question)
    if not tokens:
        return []

    matches = retrieve_resource_links(question, top_k=top_k)
    if matches:
        return matches

    for code, description in VERIFICAR_DB.items():
        searchable = f"{code} {description}"
        normalized = normalize_search_text(searchable)
        score = sum(4 if token in normalize_search_text(code) else 1 for token in tokens if token in normalized)
        if score:
            matches.append(
                {
                    "source": "knowledge_base/verificar/db_codigos.js",
                    "chunk": 1,
                    "text": f"Codigo {code}: {description}",
                    "score": float(score),
                    "kind": "verificar_db",
                    "code": code,
                    "description": description,
                }
            )

    if matches:
        return sorted(matches, key=lambda item: item["score"], reverse=True)[:top_k]

    for document in load_local_documents():
        for index, chunk in enumerate(chunk_text(document.text, chunk_size=900, overlap=100), start=1):
            normalized = normalize_search_text(chunk)
            score = sum(1 for token in tokens if token in normalized)
            if score:
                matches.append(
                    {
                        "source": document.source,
                        "chunk": index,
                        "text": chunk,
                        "score": float(score),
                        "kind": "local_document",
                    }
                )

    return sorted(matches, key=lambda item: item["score"], reverse=True)[:top_k]


def generate_local_answer(question: str, matches: list[dict]) -> str:
    if not matches:
        return "No encuentro esa informacion en mi base de conocimiento."

    link_matches = [match for match in matches if match.get("kind") == "resource_link"]
    if link_matches:
        lines = ["Encontré estos enlaces listos para abrir:", ""]
        lines.append("| Tipo | Recurso | Enlace |")
        lines.append("|---|---|---|")
        seen_urls = set()
        for match in link_matches[:6]:
            url = str(match.get("url", ""))
            if url in seen_urls:
                continue
            seen_urls.add(url)
            kind = str(match.get("resource_kind", "")).replace("|", "/") or "link"
            title = str(match.get("title", "")).replace("|", "/") or url
            lines.append(f"| {kind} | {title} | {url} |")
        return "\n".join(lines)

    verificar_matches = [match for match in matches if match.get("kind") == "verificar_db"]
    if verificar_matches:
        lines = ["Encontré esta información en la base de VERIFICAR:", ""]
        lines.append("| Codigo | Descripcion |")
        lines.append("|---|---|")
        seen_codes = set()
        for match in verificar_matches[:8]:
            code = str(match.get("code", ""))
            if code in seen_codes:
                continue
            seen_codes.add(code)
            description = str(match.get("description", "")).replace("|", "/")
            lines.append(f"| {code} | {description} |")
        if "muestr" in normalize_search_text(question) and not any(Path(str(match.get("source", ""))).suffix.lower() in IMAGE_FILE_TYPES for match in matches):
            lines.append("")
            lines.append("No tengo una imagen asociada a esa busqueda en la base; solo encontre la referencia textual.")
        return "\n".join(lines)

    best = matches[0]
    snippet = " ".join(str(best.get("text", "")).split())
    if len(snippet) > 700:
        snippet = snippet[:700].rsplit(" ", 1)[0] + "..."
    return f"Encontré esto en la base de conocimiento:\n\n{snippet}"


def answer_with_local_knowledge(question: str, top_k: int = 8, history: list[dict] | None = None) -> dict:
    search_question = contextual_question(question, history)
    if wants_external_link(question):
        link_matches = retrieve_resource_links(search_question, top_k=top_k)
        if link_matches:
            return {"answer": generate_local_answer(question, link_matches), "sources": link_matches}
        return {
            "answer": "No tengo un enlace abierto para esa consulta todavia. Puedes agregarlo en `knowledge_base/links/` y lo entregare directamente la proxima vez.",
            "sources": [],
        }

    matches = retrieve_local(search_question, top_k=top_k)
    return {"answer": generate_local_answer(question, matches), "sources": matches}


def generate_answer(question: str, matches: list[dict], user_name: str = "", history: list[dict] | None = None) -> str:
    # BLOQUE 5: respuesta con RAG/OpenAI cuando existe indice y API key.
    load_index()
    client = get_client()
    model = os.getenv("OPENAI_MODEL", "gpt-5.4-mini")
    context = build_context(matches)
    try:
        response = client.responses.create(
            model=model,
            input=[
                {
                    "role": "system",
                    "content": (
                        "Eres un asistente interno de consulta general para la empresa. "
                        "Tu funcion es ayudar a los colaboradores a encontrar informacion en la base de conocimiento "
                        "sobre codigos, acabados, colores, cronogramas, procesos, enlaces, documentos y material operativo. "
                        "Responde de forma clara, cercana y amable, como un companero experto que ayuda a encontrar rapido la informacion. "
                        "Empieza con una respuesta corta y util; luego agrega detalles si ayudan. "
                        "Cuando la respuesta incluya listas de datos, comparaciones, codigos, lotes, fechas, enlaces, colores, responsables o estados, "
                        "organiza la informacion en una tabla Markdown para que sea facil de leer. "
                        "Usa bullets solo cuando una tabla no aporte claridad. "
                        "Usa unicamente informacion del contexto entregado; no inventes datos. "
                        "Usa el historial de conversacion solo para entender continuidad, referencias y pronombres como "
                        "'eso', 'ese', 'el anterior' o 'lo mismo'. No uses el historial para inventar informacion que no este respaldada por el contexto. "
                        "Si la pregunta es ambigua, responde con la informacion mas cercana y sugiere como precisar la consulta. "
                        "Si el contexto no contiene la respuesta, di exactamente: "
                        "'No encuentro esa informacion en mi base de conocimiento.' "
                        "No muestres fuentes, nombres de archivo, fragmentos ni citas en la respuesta visible."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Usuario: {user_name or 'Usuario interno'}\n\n"
                        f"Historial reciente:\n{history_text(history) or 'Sin historial reciente.'}\n\n"
                        f"Contexto:\n{context}\n\n"
                        f"Pregunta: {question}"
                    ),
                },
            ],
        )
    except Exception as error:
        raise friendly_openai_error(error) from error
    return strip_visible_sources(response.output_text)


def answer_question(question: str, top_k: int = 5) -> str:
    if wants_external_link(question):
        return answer_with_local_knowledge(question, top_k=top_k)["answer"]

    verificar_result = answer_verificar_question(question)
    if verificar_result:
        return verificar_result["answer"]

    try:
        matches = retrieve(question, top_k=top_k)
        return generate_answer(question, matches)
    except UserFacingError:
        return answer_with_local_knowledge(question, top_k=top_k)["answer"]


def answer_question_with_sources(question: str, top_k: int = 5, user_name: str = "", history: list[dict] | None = None) -> dict:
    if wants_external_link(question):
        return answer_with_local_knowledge(question, top_k=top_k, history=history)

    contextual = contextual_question(question, history)
    verificar_result = answer_verificar_question(contextual)
    if verificar_result:
        return verificar_result

    try:
        matches = retrieve(contextual, top_k=top_k)
        return {"answer": generate_answer(question, matches, user_name=user_name, history=history), "sources": matches}
    except UserFacingError:
        return answer_with_local_knowledge(question, top_k=top_k, history=history)
